from pptx import Presentation
from pptx.util import Inches

def generate_presentation():
    """
    Generates a PowerPoint presentation with fleet data.
    """
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Fleet Management Report"
    subtitle.text = "Quarterly Review"

    # Slide 2: Executive Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Executive Summary"
    tf = body_shape.text_frame
    tf.text = "TCO increased by 5% compared to the previous quarter."
    p = tf.add_paragraph()
    p.text = "Placeholder for TCO trend chart"
    p.level = 1

    # Slide 3: EV Transition
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "EV Transition"
    tf = body_shape.text_frame
    tf.text = "Placeholder for EV transition chart"

    # Slide 4: Model Efficiency
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Model Efficiency"
    tf = body_shape.text_frame
    tf.text = "Placeholder for model efficiency chart"

    return prs
